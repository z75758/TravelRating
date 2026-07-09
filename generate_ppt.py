"""
Generate defense PPT for TravelRating system.
Uses template theme but builds all slides with fresh content.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Load template for theme/layouts
template_path = r'C:\Users\asus\OneDrive\桌面\Web开发技术课程设计-班级-学号-姓名-答辩PPT(2).pptx'
prs = Presentation(template_path)

# Remove all existing slides
xml_slides = prs.slides._sldIdLst
while len(xml_slides) > 0:
    rId = xml_slides[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    xml_slides.remove(xml_slides[0])

# ── Color constants ──
GREEN = RGBColor(0x2D, 0x6A, 0x4F)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
ACCENT_BG = RGBColor(0xE8, 0xF5, 0xE9)

# ── Helpers ──
def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False,
                 color=DARK, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_size=12,
                      color=DARK, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(4)
    return txBox

def add_section_number(slide, left, top, size, number):
    """Add a large section number."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(left), Emu(top), Emu(size), Emu(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = number
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_card_box(slide, left, top, width, height, title, lines):
    """Add a card-style box with title and content."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)
    # Add title and content as text box overlay
    add_text_box(slide, left + Emu(200000), top + Emu(120000), width - Emu(400000), Emu(400000),
                 title, font_size=16, bold=True, color=GREEN)
    add_multiline_box(slide, left + Emu(200000), top + Emu(520000), width - Emu(400000),
                      height - Emu(640000), lines, font_size=11)
    return shape

# ═══════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[0])
# Clear placeholders and set text
for shape in slide.shapes:
    if shape.has_text_frame:
        tf = shape.text_frame
        for p in tf.paragraphs:
            p.clear()
            for run in p.runs:
                run.text = ''
# Add custom title
add_text_box(slide, 1500000, 2300000, 9200000, 1200000,
             '旅游地点评选与互动系统\n设计与实现',
             font_size=36, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 3500000, 4200000, 5200000, 500000,
             '汇报时间：2026年7月8日',
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 4200000, 4800000, 3800000, 500000,
             '答辩人：张绮雯',
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════
# SLIDE 2: Table of Contents
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_text_box(slide, 2400000, 1000000, 7400000, 1000000, '目  录',
             font_size=40, bold=True, alignment=PP_ALIGN.LEFT)

toc = [
    ('01', '项目背景与需求分析'),
    ('02', '系统设计与技术架构'),
    ('03', '系统功能实现展示'),
    ('04', '总结与展望'),
]
for i, (num, title) in enumerate(toc):
    y_base = 2100000 + i * 850000
    # Number
    add_text_box(slide, 7800000, y_base, 700000, 550000, num,
                 font_size=28, bold=True, color=GREEN, alignment=PP_ALIGN.RIGHT)
    # Title
    add_text_box(slide, 3500000, y_base, 4000000, 550000, title,
                 font_size=22, color=DARK)
    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Emu(3500000), Emu(y_base + 650000),
                                  Emu(5200000), Emu(15000))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    line.line.fill.background()

# ═══════════════════════════════════════
# SLIDE 3: Section 01 divider
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[2])
add_section_number(slide, 3200000, 2400000, 1000000, '01')
add_text_box(slide, 4200000, 2550000, 5500000, 900000,
             '项目背景与需求分析',
             font_size=34, bold=True, alignment=PP_ALIGN.LEFT)

# ═══════════════════════════════════════
# SLIDE 4: Background & Requirements
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_text_box(slide, 700000, 300000, 6000000, 600000, '项目背景与需求分析',
             font_size=26, bold=True, color=DARK)

# Three column cards
card_data = [
    ('项目背景', [
        '• 旅游业蓬勃发展，游客对目的地信息获取和评价反馈需求日益增长',
        '• 传统旅游网站多采用静态展示，缺乏用户互动和评价机制',
        '• 游客出行决策缺少真实用户评分的参考依据',
        '• 旅游信息获取渠道分散，缺乏统一的评选互动平台',
        '• 目的地管理依赖人工操作，效率低下且易出错',
    ]),
    ('项目目标', [
        '• 构建集浏览、搜索、评分、评论于一体的旅游互动平台',
        '• 实现交互式五星评分机制，提供真实用户反馈数据',
        '• 提供管理员后台，实现目的地信息的高效CRUD管理',
        '• 采用H2嵌入式数据库实现零安装部署，开箱即用',
        '• 应用Taste-Skill极简设计系统，打造专业美观界面',
    ]),
    ('系统特色', [
        '• 交互式五星评分：hover预览+click提交+AJAX异步',
        '• 每用户每目的地仅一票：UNIQUE约束+UPSERT策略',
        '• Taste-Skill设计：暖白色调+1px边框+8px圆角',
        '• H2嵌入式数据库：零配置，MySQL兼容语法',
        '• SHA-256+Salt密码加密 + Filter双层鉴权',
        '• 响应式布局 + 暗色模式自适应',
    ]),
]

for i, (title, lines) in enumerate(card_data):
    left = 400000 + i * 3800000
    add_card_box(slide, left, 1100000, 3600000, 5400000, title, lines)

# ═══════════════════════════════════════
# SLIDE 5: Section 02 divider
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[2])
add_section_number(slide, 3200000, 2400000, 1000000, '02')
add_text_box(slide, 4200000, 2550000, 5500000, 900000,
             '系统设计与技术架构',
             font_size=34, bold=True, alignment=PP_ALIGN.LEFT)

# ═══════════════════════════════════════
# SLIDE 6: Tech Architecture
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_text_box(slide, 700000, 300000, 6000000, 600000, '系统设计与技术架构',
             font_size=26, bold=True, color=DARK)

# Left: Tech stack
add_text_box(slide, 400000, 1100000, 5400000, 400000, '技术栈',
             font_size=20, bold=True, color=GREEN)
add_multiline_box(slide, 400000, 1500000, 5400000, 2400000, [
    '【后端】',
    '• 语言：Java 11',
    '• 框架：JSP + Servlet + JDBC（MVC模式）',
    '• 数据库：H2 Database 2.3.232（MySQL兼容）',
    '• 服务器：Apache Tomcat 9.0.97（嵌入式）',
    '• 构建：Maven 3.9 + tomcat7-maven-plugin',
    '',
    '【前端】',
    '• 视图：JSP（16个页面） + JSTL + EL表达式',
    '• 样式：Taste-Skill CSS设计系统（全局变量）',
    '• 脚本：Vanilla JavaScript（无框架依赖）',
    '• 异步：Fetch API + AJAX（评分/上传）',
    '• 字体：Inter Tight + Newsreader（Google Fonts）',
], font_size=11)

# Right: Architecture
add_text_box(slide, 6200000, 1100000, 5400000, 400000, 'MVC分层架构',
             font_size=20, bold=True, color=GREEN)
add_multiline_box(slide, 6200000, 1500000, 5400000, 5000000, [
    '┌─ View 层 ─────────────────┐',
    '│ 16个JSP + 2个CSS + JS脚本  │',
    '│ Taste-Skill 设计令牌       │',
    '│ 响应式 + 暗色模式          │',
    '└───────────────────────────┘',
    '          ↕',
    '┌─ Controller 层 ───────────┐',
    '│ 6个Servlet（@WebServlet）  │',
    '│ 路由映射 + 权限验证        │',
    '│ PRG模式 / JSON API         │',
    '└───────────────────────────┘',
    '          ↕',
    '┌─ Service 层 ──────────────┐',
    '│ 4个Service（业务规则）     │',
    '│ 评分校验 + 密码加密        │',
    '│ 投票去重 + 评分同步        │',
    '└───────────────────────────┘',
    '          ↕',
    '┌─ DAO 层 ──────────────────┐',
    '│ 4个DAO（纯JDBC）           │',
    '│ PreparedStatement 防注入   │',
    '│ try-with-resources 管理    │',
    '└───────────────────────────┘',
    '          ↕',
    '┌─ Model 层 ────────────────┐',
    '│ User / Destination         │',
    '│ Comment / Vote             │',
    '└───────────────────────────┘',
], font_size=10, color=RGBColor(0x33, 0x33, 0x33))

# ═══════════════════════════════════════
# SLIDE 7: Database Design
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_text_box(slide, 700000, 300000, 6000000, 600000, '数据库设计',
             font_size=26, bold=True, color=DARK)

# 4 tables as cards
db_cards = [
    ('users（用户表）', '主键: id\n唯一: username, email\n字段: password(SHA-256+Salt), avatar, role(user/admin), created_at, updated_at\n约束: username UNIQUE, email UNIQUE'),
    ('destinations（目的地表）', '主键: id\n字段: name, region, type, image, description, address, open_time, ticket_price(DECIMAL), rating(派生), popularity, created_by(FK→users)\n约束: FK ON DELETE SET NULL'),
    ('comments（评论表）', '主键: id\n字段: destination_id(FK), user_id(FK), content(TEXT), rating(1-5), created_at\n约束: FK→destinations ON DELETE CASCADE, FK→users ON DELETE CASCADE'),
    ('votes（投票表）', '主键: id\n字段: destination_id(FK), user_id(FK), score(1-5), created_at\n约束: UNIQUE(destination_id, user_id) ★每用户每目的地仅一票, FK ON DELETE CASCADE'),
]

for i, (title, desc) in enumerate(db_cards):
    x = 300000 + (i % 2) * 5900000
    y = 1100000 + (i // 2) * 2700000
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Emu(x), Emu(y), Emu(5600000), Emu(2500000))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    card.line.width = Pt(1)
    add_text_box(slide, x + 200000, y + 100000, 5200000, 400000,
                 title, font_size=16, bold=True, color=GREEN)
    add_multiline_box(slide, x + 200000, y + 500000, 5200000, 1900000,
                      desc.split('\n'), font_size=10)

# ═══════════════════════════════════════
# SLIDE 8: Innovation & Features
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_text_box(slide, 700000, 300000, 6000000, 600000, '系统创新性与项目价值',
             font_size=26, bold=True, color=DARK)

innovations = [
    ('01', '交互式五星评分', [
        'hover预览：鼠标悬停高亮变金黄色并显示文字提示',
        'click提交：Fetch API异步POST，无需刷新页面',
        'UPSERT策略：已评分则UPDATE，未评分则INSERT',
        '实时更新：评分后立即刷新平均分、星级、投票数',
        'UNIQUE约束：每用户每目的地仅一票',
    ]),
    ('02', '安全防护体系', [
        'SHA-256 + 16字节随机盐值加密密码',
        'AuthFilter双层鉴权（登录状态 + 管理员角色）',
        'PreparedStatement参数化查询全面防SQL注入',
        'JSP自动转义 + Filter链防护XSS攻击',
        'Session 30分钟超时 + 安全退出',
    ]),
    ('03', '设计系统与体验', [
        'Taste-Skill CSS自定义属性统一视觉语言',
        'Bento Grid 12列非对称网格首页布局',
        'IntersectionObserver滚动渐入动画',
        '响应式布局768px断点 + 暗色模式自适应',
        'Toast消息通知 + 按钮scale(0.98)反馈',
    ]),
]

for i, (num, title, lines) in enumerate(innovations):
    x = 300000 + i * 3900000
    add_text_box(slide, x, 1100000, 3700000, 400000,
                 f'{num}  {title}', font_size=16, bold=True, color=GREEN)
    add_multiline_box(slide, x, 1500000, 3700000, 3500000, lines, font_size=11)

# Value summary at bottom
add_multiline_box(slide, 300000, 5200000, 11400000, 1400000, [
    '项目价值：为游客提供便捷的旅游信息获取和社区评价互动平台  |  通过用户评分数据聚合形成社区驱动的目的地推荐  |  为旅游数字化转型提供实践案例',
], font_size=13, color=GRAY)

# ═══════════════════════════════════════
# SLIDE 9: Section 03 divider
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[2])
add_section_number(slide, 3200000, 2400000, 1000000, '03')
add_text_box(slide, 4200000, 2550000, 5500000, 900000,
             '系统功能实现展示',
             font_size=34, bold=True, alignment=PP_ALIGN.LEFT)

# ═══════════════════════════════════════
# SLIDE 10: Frontend Pages Showcase
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_text_box(slide, 700000, 250000, 8000000, 500000, '前台核心页面展示',
             font_size=24, bold=True, color=DARK)

frontend_pages = [
    ('系统首页', [
        '• Bento Grid 12列非对称网格',
        '• 6大热门目的地卡片展示',
        '• 三合一搜索栏（关键字+地区+类型）',
        '• 4列统计卡片（目的地/地区/类型/投票数）',
        '• Hero区域 + CTA按钮引导',
    ]),
    ('目的地列表页', [
        '• 表格视图展示全部12个目的地',
        '• 筛选工具栏（地区+类型+排序）',
        '• 关键字模糊搜索（名称+描述）',
        '• 评分/热度/最新三种排序',
        '• 卡片悬停上浮动效',
    ]),
    ('目的地详情页', [
        '• 大幅封面图（1200×600）',
        '• 信息侧边栏（地址/时间/票价/热度）',
        '• ★交互式五星评分组件',
        '• 评论区（发表+列表+删除）',
        '• Toast消息通知反馈',
    ]),
    ('登录 & 注册页', [
        '• 中央居中卡片布局',
        '• 支持用户名或邮箱登录',
        '• 前端密码一致性校验',
        '• 后端用户名/邮箱唯一性检查',
        '• SHA-256+Salt加密存储',
    ]),
]

for i, (title, lines) in enumerate(frontend_pages):
    x = 300000 + (i % 2) * 5900000
    y = 900000 + (i // 2) * 2800000
    add_card_box(slide, x, y, 5600000, 2600000, title, lines)

# ═══════════════════════════════════════
# SLIDE 11: Core Features - Rating & Comment
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_text_box(slide, 700000, 250000, 8000000, 500000, '核心互动功能：星级评分与评论管理',
             font_size=24, bold=True, color=DARK)

# Left: Rating mechanism
add_text_box(slide, 400000, 900000, 5400000, 400000,
             '⭐ 五星评分机制', font_size=20, bold=True, color=GREEN)
add_multiline_box(slide, 400000, 1300000, 5400000, 5000000, [
    '【前端交互流程】',
    '1. hover → 星星变金黄(#e6a817) + 文字提示(很差/较差/一般/较好/非常好)',
    '2. click → 构建URLSearchParams → Fetch POST /vote',
    '3. 解析JSON → 更新currentScore → highlight() → Toast通知',
    '4. 防重复提交：isSubmitting布尔锁',
    '',
    '【后端处理流程】',
    '1. VoteServlet接收AJAX请求 → 验证登录(401 if not)',
    '2. VoteService.rate() → clamp评分1-5',
    '3. VoteDAO.rate() → findByUserAndDestination()',
    '   存在？→ UPDATE score   不存在？→ INSERT',
    '4. DestinationDAO.updateRating() → SELECT AVG(score)',
    '5. 返回JSON: {success, avgRating, voteCount, userScore}',
    '',
    '【API端点】',
    '• POST /vote?action=rate      → 提交/修改评分',
    '• POST /vote?action=unrate    → 取消评分',
    '• GET  /vote?destinationId=X  → 查询状态',
], font_size=11)

# Right: Comment management
add_text_box(slide, 6200000, 900000, 5400000, 400000,
             '💬 评论管理功能', font_size=20, bold=True, color=GREEN)
add_multiline_box(slide, 6200000, 1300000, 5400000, 5000000, [
    '【发表评论】',
    '1. 登录用户在详情页输入文字内容',
    '2. 表单POST /comment（accept-charset="UTF-8"）',
    '3. CommentServlet → CommentService.addComment()',
    '4. CommentDAO.create() → INSERT INTO comments',
    '5. 302重定向 → 页面刷新 → 新评论显示在列表顶部',
    '',
    '【删除评论】★新增功能',
    '1. 每条评论右侧显示"×"删除按钮',
    '2. 按钮仅对评论作者和管理员可见',
    '3. 点击 → confirm("确定要删除这条评论吗？")',
    '4. POST /comment?action=delete',
    '5. CommentService.deleteComment() 三重校验：',
    '   ① 评论存在性（findById）',
    '   ② 作者匹配 或 管理员越权',
    '   ③ 通过后 → DAO.delete() 物理删除',
    '',
    '【中文乱码修复】★关键修复',
    '• 移除@WebFilter双重重注册',
    '• web.xml显式声明Filter执行顺序（Encoding先于Auth）',
    '• 防御性request.setCharacterEncoding("UTF-8")',
], font_size=11)

# ═══════════════════════════════════════
# SLIDE 12: Admin Features
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_text_box(slide, 700000, 250000, 8000000, 500000, '管理后台功能展示',
             font_size=24, bold=True, color=DARK)

admin_features = [
    ('后台仪表盘', [
        '• 4个统计卡片：目的地/用户/投票总数',
        '• 管理员专属入口（AuthFilter拦截）',
        '• 登录后自动跳转仪表盘',
        '• 全局导航栏显示"Admin"入口链接',
    ]),
    ('目的地CRUD管理', [
        '• 创建：8字段表单 + AJAX图片上传',
        '• 编辑：预填充表单 + UPDATE提交',
        '• 删除：confirm确认 + 级联删除评论投票',
        '• 图片上传：File API预览 + FormData POST',
        '• UUID文件名 + 类型校验(JPG/PNG/GIF/WebP)',
        '• 大小限制：最大5MB',
    ]),
    ('用户管理', [
        '• 用户列表查看（ID/用户名/邮箱/角色/时间）',
        '• 管理员可删除用户',
        '• 普通用户不可访问管理路径（302重定向）',
        '• 角色区分：user vs admin',
    ]),
    ('搜索与数据操作', [
        '• 多条件组合搜索（关键字+地区+类型+排序）',
        '• INSERT：PreparedStatement参数化',
        '• UPDATE：仅更新可变字段',
        '• DELETE：外键ON DELETE CASCADE级联',
        '• AVG聚合：实时计算平均评分',
    ]),
]

for i, (title, lines) in enumerate(admin_features):
    x = 300000 + (i % 2) * 5900000
    y = 900000 + (i // 2) * 2800000
    add_card_box(slide, x, y, 5600000, 2600000, title, lines)

# ═══════════════════════════════════════
# SLIDE 13: Search & Personal Center
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_text_box(slide, 700000, 250000, 8000000, 500000, '搜索筛选与个人中心',
             font_size=24, bold=True, color=DARK)

add_card_box(slide, 300000, 900000, 5600000, 2600000, '多条件搜索与筛选', [
    '• 关键字模糊匹配：name LIKE ? OR description LIKE ?',
    '• 地区精确筛选：WHERE region = ?',
    '• 类型精确筛选：WHERE type = ?',
    '• 三合一首页搜索栏（关键字+地区下拉+类型下拉）',
    '• StringBuilder动态拼接WHERE条件',
    '• 三种排序：popularity DESC / rating DESC / created_at DESC',
    '• 支持任意条件组合（全部可选）',
])

add_card_box(slide, 6200000, 900000, 5600000, 2600000, '个人中心（user/profile.jsp）', [
    '• 资料编辑：修改用户名/邮箱/头像URL',
    '• 密码修改：验证原密码 + 新密码SHA-256加密',
    '• Session同步：更新后刷新session中User对象',
    '• 表单预填充当前信息',
    '• 前端校验 + 后端双重验证',
])

add_card_box(slide, 300000, 3700000, 5600000, 2600000, 'Taste-Skill设计系统应用', [
    '• CSS自定义属性：全局颜色/间距/圆角/阴影令牌',
    '• 暖白色(#FBFBFA)背景 + 1px细边框(#EAEAEA)',
    '• 8px统一圆角 + 微妙阴影(rgba 0.04)',
    '• IntersectionObserver滚动渐入动画',
    '• prefers-color-scheme暗色模式自动切换',
    '• prefers-reduced-motion尊重用户偏好',
    '• 768px响应式断点 + 触控友好(44px+)',
])

add_card_box(slide, 6200000, 3700000, 5600000, 2600000, '安全性设计总结', [
    '• 密码：SHA-256 + 16字节随机盐（格式："salt:hash"）',
    '• SQL注入：全PreparedStatement参数化查询',
    '• XSS防护：JSP自动转义HTML实体',
    '• 认证：AuthFilter拦截所有非公开路径',
    '• 授权：user.isAdmin()双重校验管理员权限',
    '• Session：30分钟超时 + session.invalidate()退出',
    '• 编码：EncodingFilter → web.xml显式顺序 → UTF-8',
    '• 文件上传：扩展名白名单 + 大小5MB限制',
])

# ═══════════════════════════════════════
# SLIDE 14: Technical Highlights
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_text_box(slide, 700000, 250000, 8000000, 500000, '关键技术实现细节',
             font_size=24, bold=True, color=DARK)

add_card_box(slide, 300000, 900000, 5600000, 2200000, 'DBUtil：双重检查锁定单例初始化', [
    '• volatile + synchronized (Double-Checked Locking)',
    '• 首次getConnection()时自动建表+插入12条种子数据',
    '• H2 URL: jdbc:h2:file:./data/travel_db;MODE=MySQL',
    '• DB_CLOSE_DELAY=-1：JVM退出前保持连接',
    '• DATABASE_TO_LOWER=TRUE：统一小写表名',
    '• initDatabase()：DROP→CREATE→INSERT，幂等初始化',
    '• close(AutoCloseable...)：可变参数静默关闭工具方法',
])

add_card_box(slide, 6200000, 900000, 5600000, 2200000, 'Filter链：编码 → 认证 → Servlet', [
    '• web.xml显式声明Filter执行顺序（解决@WebFilter无序问题）',
    '• EncodingFilter：request.setCharacterEncoding("UTF-8")',
    '• AuthFilter三级检查：静态资源放行 → 公开路径放行 → Session验证',
    '• 管理员路径额外校验：user.isAdmin() → 否则302到首页',
    '• 白名单路径：css/js/images/uploads + 登录/注册/首页',
    '• 未登录访问受保护页 → 302重定向到login.jsp',
])

add_card_box(slide, 300000, 3300000, 5600000, 2200000, 'VoteDAO：UPSERT评分策略', [
    '• rate(destId, userId, score)：先查后改',
    '• 存在 → UPDATE votes SET score=? WHERE id=?',
    '• 不存在 → INSERT INTO votes (...) VALUES (...)',
    '• getAverageScore(destId)：SELECT COALESCE(AVG(score),0)',
    '• removeRating(destId, userId)：DELETE + 重算AVG',
    '• 精度：Math.round(avg * 10.0) / 10.0（1位小数）',
])

add_card_box(slide, 6200000, 3300000, 5600000, 2200000, '项目统计', [
    '• 总代码行数：约3,500行（Java + JSP + CSS + JS）',
    '• Servlet控制器：6个 | DAO数据访问：4个',
    '• Service业务逻辑：4个 | Model实体：4个',
    '• Filter过滤器：2个 | Util工具类：3个',
    '• JSP页面：16个 | CSS样式文件：2个',
    '• 数据库表：4张 | 示例数据：12条目的地',
    '• 测试用例：97项全部通过（100%通过率）',
])

# ═══════════════════════════════════════
# SLIDE 15: Section 04 divider
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[2])
add_section_number(slide, 3200000, 2400000, 1000000, '04')
add_text_box(slide, 4200000, 2550000, 5500000, 900000,
             '总结与展望',
             font_size=34, bold=True, alignment=PP_ALIGN.LEFT)

# ═══════════════════════════════════════
# SLIDE 16: Summary
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_text_box(slide, 700000, 250000, 8000000, 500000, '总结与展望',
             font_size=26, bold=True, color=DARK)

# Problems solved
add_text_box(slide, 400000, 900000, 5400000, 400000, '项目开发中解决的关键问题',
             font_size=18, bold=True, color=GREEN)
add_multiline_box(slide, 400000, 1300000, 5400000, 5000000, [
    '1. 中文乱码修复',
    '   移除@WebFilter双重重注册，在web.xml中显式声明Filter执行顺序（EncodingFilter先于AuthFilter），',
    '   配合CommentServlet防御性setCharacterEncoding调用，彻底解决POST请求中文参数乱码。',
    '',
    '2. 评分UPSERT逻辑',
    '   通过数据库UNIQUE(destination_id, user_id)联合约束 + 应用层findByUserAndDestination →',
    '   UPDATE/INSERT分支策略，实现优雅的评分创建/修改二合一，避免唯一约束冲突。',
    '',
    '3. Servlet 3.0兼容适配',
    '   将web.xml命名空间从Servlet 3.1降级至3.0以兼容Tomcat 7嵌入式容器（tomcat7-maven-plugin 2.2），',
    '   同时保留metadata-complete="false"启用@WebServlet/@WebFilter注解扫描。',
    '',
    '4. 评论删除权限控制',
    '   通过CommentService.deleteComment()实现三层校验：',
    '   ① findById检查评论存在性 → ② 作者ID匹配或isAdmin()校验 → ③ DAO.delete()物理删除。',
    '   前端JSP条件渲染删除按钮（仅作者或管理员可见），JavaScript confirm()二次确认防止误删。',
], font_size=10)

# Future work
add_text_box(slide, 6200000, 900000, 5400000, 400000, '未来改进方向',
             font_size=18, bold=True, color=GREEN)
add_multiline_box(slide, 6200000, 1300000, 5400000, 5000000, [
    '1. 密码加密升级',
    '   SHA-256 → bcrypt/scrypt/PBKDF2慢哈希算法',
    '',
    '2. 数据持久化',
    '   DBUtil中DROP TABLE改为CREATE TABLE IF NOT EXISTS，保留用户数据',
    '',
    '3. CSRF与安全加固',
    '   为表单添加CSRF Token + 登录失败次数限制',
    '',
    '4. 日志与监控',
    '   printStackTrace → SLF4J/Logback统一日志框架',
    '',
    '5. 单元测试',
    '   JUnit 5 + Mockito + H2内存库自动化测试',
    '',
    '6. 移动端扩展',
    '   Flutter跨平台App开发 + RESTful API重构',
    '',
    '7. 智能推荐',
    '   基于协同过滤算法的个性化目的地推荐',
    '',
    '8. 在线支付与预订',
    '   微信/支付宝集成 + 在线票务预订系统',
], font_size=10)

# ═══════════════════════════════════════
# SLIDE 17: Thank You
# ═══════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[10])
add_text_box(slide, 1500000, 2200000, 9200000, 1200000,
             '敬请老师批评指正',
             font_size=40, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 4000000, 4000000, 4200000, 500000,
             '答辩人：张绮雯', font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 3800000, 4500000, 4600000, 500000,
             '汇报时间：2026年7月8日', font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# ── Save ──
output = r'F:\MyProject\TravelRating\TravelRating_答辩PPT.pptx'
prs.save(output)
print(f'PPT saved: {output}')
print(f'Size: {os.path.getsize(output)/1024:.0f} KB | Slides: {len(prs.slides)}')
